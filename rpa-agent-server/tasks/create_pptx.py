"""
create_pptx.py — 使用 python-pptx 生成真实 .pptx 文件

输入 taskParams:
  deliverableId : str        — 对应 LTC-OS deliverables 表的 ID（用于文件命名）
  title         : str        — 方案标题（也是文件名前缀）
  slides        : list[dict] — 幻灯片内容
    每个 slide:
      title   : str        — 幻灯片标题
      bullets : list[str]  — 要点列表（优先使用，AI 输出格式）
      content : str (可选) — 纯文本内容（换行分割，兼容旧格式）
      notes   : str (可选) — 备注
      accent  : str (可选) — 右侧数字卡片（如 "85%"、"300人"）
  style         : str (可选) — "tech"（默认）| "business" | "custom:accentHex,bgHex"
  companyName   : str (可选) — 副标题/公司名

输出:
  { fileUrl: "/files/xxx.pptx" }
"""

import os
import uuid
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


STORAGE_DIR = os.path.join(os.path.dirname(__file__), "..", "storage")

# ─── 主题色定义 ───────────────────────────────────────────────────────────────

# 科技感深色主题（默认）
TECH_THEME = {
    "bg":      RGBColor(0x0D, 0x1B, 0x2A),   # 近黑深蓝底
    "primary": RGBColor(0x0D, 0x1B, 0x2A),   # 同底色（内容页背景）
    "accent":  RGBColor(0x00, 0xD4, 0xFF),   # 亮青色
    "accent2": RGBColor(0x7B, 0x2F, 0xFF),   # 紫色辅助
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "text":    RGBColor(0xE8, 0xF4, 0xFF),   # 浅蓝白
    "subtext": RGBColor(0x9B, 0xC4, 0xE2),   # 蓝灰副文字
    "header_bar": RGBColor(0x00, 0xD4, 0xFF), # 标题条
    "header_text": RGBColor(0x0D, 0x1B, 0x2A), # 标题条文字
}

# 商务简洁浅色主题
BUSINESS_THEME = {
    "bg":      RGBColor(0xFF, 0xFF, 0xFF),
    "primary": RGBColor(0x1A, 0x3A, 0x6B),
    "accent":  RGBColor(0x00, 0x8B, 0xD8),
    "accent2": RGBColor(0xFF, 0x6B, 0x35),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "text":    RGBColor(0x33, 0x33, 0x33),
    "subtext": RGBColor(0x66, 0x66, 0x66),
    "header_bar": RGBColor(0x1A, 0x3A, 0x6B),
    "header_text": RGBColor(0xFF, 0xFF, 0xFF),
}


def _parse_hex(hex_str: str) -> RGBColor:
    """解析6位十六进制颜色字符串（不含#）"""
    try:
        h = hex_str.strip().lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(0x00, 0xD4, 0xFF)


def _get_theme(style: str) -> dict:
    """根据 style 参数返回主题色字典"""
    if style == "business":
        return BUSINESS_THEME
    if style.startswith("custom:"):
        parts = style[7:].split(",")
        accent = _parse_hex(parts[0]) if len(parts) > 0 and parts[0] else TECH_THEME["accent"]
        bg_color = _parse_hex(parts[1]) if len(parts) > 1 and parts[1] else TECH_THEME["bg"]
        theme = dict(TECH_THEME)
        theme["accent"] = accent
        theme["bg"] = bg_color
        theme["primary"] = bg_color
        theme["header_bar"] = accent
        # 深色背景用深色文字（亮青），浅色背景用深色文字
        lum = (bg_color[0] * 299 + bg_color[1] * 587 + bg_color[2] * 114) // 1000
        if lum > 128:
            theme["text"] = RGBColor(0x22, 0x22, 0x22)
            theme["header_text"] = RGBColor(0xFF, 0xFF, 0xFF)
        return theme
    # 默认 tech
    return TECH_THEME


def _set_slide_bg(slide, color: RGBColor):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_cover_slide(prs: Presentation, title: str, subtitle: str = "", theme: dict = None):
    """封面：深色背景 + 亮色左侧粗竖条 + 白色大标题"""
    if theme is None:
        theme = TECH_THEME

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme["bg"])

    # 左侧粗竖条（亮色）
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.45), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

    # 装饰：右侧大半透明圆形（用细框矩形模拟）
    deco = slide.shapes.add_shape(
        9,  # oval
        Inches(9.5), Inches(1.5), Inches(4.8), Inches(4.8)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = theme["accent2"] if "accent2" in theme else theme["accent"]
    deco.line.fill.background()
    # 设置透明度（通过 XML 直接操作）
    try:
        from pptx.oxml.ns import qn
        solidFill = deco.fill._xPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                from lxml import etree
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', '12000')  # 12% 不透明（88% 透明）
    except Exception:
        pass

    # 主标题
    txBox = slide.shapes.add_textbox(Inches(0.65), Inches(2.2), Inches(8.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = theme["white"]

    # 分隔线
    line_box = slide.shapes.add_shape(1, Inches(0.65), Inches(4.15), Inches(5.0), Inches(0.04))
    line_box.fill.solid()
    line_box.fill.fore_color.rgb = theme["accent"]
    line_box.line.fill.background()

    # 副标题 / 公司名
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.65), Inches(4.35), Inches(8.5), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = theme["accent"]

    return slide


def _add_content_slide(
    prs: Presentation,
    slide_title: str,
    bullets: list[str],
    notes: str = "",
    accent_text: str = "",
    theme: dict = None,
):
    """内容页：深色背景 + 亮色标题条 + bullet 列表（含 ▶ 前缀）
    - ≤4 bullets: 单栏，字号 Pt(20)
    - ≥5 bullets: 双栏，字号 Pt(17)
    accent_text: 右侧数字卡片（如 "85%"）
    """
    if theme is None:
        theme = TECH_THEME

    is_dark = _is_dark_theme(theme)
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme["bg"])

    # 顶部标题条
    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = theme["header_bar"]
    header_bar.line.fill.background()

    # 左侧小色条（深色内容区用 accent2）
    left_bar_color = theme.get("accent2", theme["accent"]) if is_dark else theme["primary"]
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.2), Inches(1.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = left_bar_color
    accent_bar.line.fill.background()

    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(12.5), Inches(0.85))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = theme["header_text"]

    # ── 布局决策 ─────────────────────────────────────────────────────────────
    two_column = len(bullets) >= 5
    has_accent_card = bool(accent_text)

    if two_column:
        _add_two_column_bullets(slide, bullets, theme, is_dark)
    else:
        # 单栏：根据是否有 accent 卡片决定宽度
        content_width = Inches(9.8) if has_accent_card else Inches(12.3)
        _add_single_column_bullets(slide, bullets, theme, is_dark, content_width)

    # accent 数字卡片（右下角）
    if has_accent_card and not two_column:
        _add_accent_card(slide, accent_text, theme)

    # 备注
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def _is_dark_theme(theme: dict) -> bool:
    """判断是否为深色主题（背景亮度 < 128 视为深色）"""
    bg = theme["bg"]
    lum = (bg[0] * 299 + bg[1] * 587 + bg[2] * 114) // 1000
    return lum < 128


def _add_single_column_bullets(slide, bullets: list[str], theme: dict, is_dark: bool, content_width):
    """单栏 bullet 列表"""
    bullet_count = len(bullets)
    if bullet_count <= 2:
        font_size = Pt(22)
    elif bullet_count <= 4:
        font_size = Pt(19)
    else:
        font_size = Pt(16)

    text_color = theme["text"] if is_dark else theme["text"]
    accent_color = theme["accent"]

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), content_width, Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_bullet_paragraph(para, bullet, font_size, text_color, accent_color)
        para.space_before = Pt(10)
        para.space_after = Pt(4)


def _add_two_column_bullets(slide, bullets: list[str], theme: dict, is_dark: bool):
    """双栏 bullet 列表（≥5 bullets 时使用）"""
    mid = (len(bullets) + 1) // 2
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    text_color = theme["text"] if is_dark else theme["text"]
    accent_color = theme["accent"]
    font_size = Pt(17)

    # 左栏
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, b in enumerate(left_bullets):
        para = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        _set_bullet_paragraph(para, b, font_size, text_color, accent_color)
        para.space_before = Pt(8)
        para.space_after = Pt(3)

    # 右栏背景框
    right_bg = slide.shapes.add_shape(1, Inches(6.85), Inches(1.2), Inches(6.0), Inches(6.0))
    right_bg.fill.solid()
    right_bg_color = RGBColor(
        max(0, theme["bg"][0] + 20),
        max(0, theme["bg"][1] + 20),
        max(0, theme["bg"][2] + 20),
    ) if is_dark else RGBColor(0xF0, 0xF5, 0xFF)
    right_bg.fill.fore_color.rgb = right_bg_color
    right_bg.line.fill.background()

    # 右栏文字
    right_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.35), Inches(5.5), Inches(5.7))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, b in enumerate(right_bullets):
        para = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
        _set_bullet_paragraph(para, b, font_size, text_color, accent_color)
        para.space_before = Pt(8)
        para.space_after = Pt(3)


def _set_bullet_paragraph(para, bullet_text: str, font_size, text_color: RGBColor, accent_color: RGBColor):
    """设置 bullet 段落：▶ 前缀用亮色，正文用 text_color"""
    from pptx.oxml.ns import qn
    from lxml import etree

    # 清除段落（若有旧内容）
    for run in list(para.runs):
        para._p.remove(run._r)

    # run1: ▶ 前缀（亮色）
    r1 = para.add_run()
    r1.text = "▶ "
    r1.font.size = font_size
    r1.font.bold = True
    r1.font.color.rgb = accent_color

    # run2: bullet 正文
    clean = bullet_text.lstrip("•-▸◦▶ \t")
    r2 = para.add_run()
    r2.text = clean
    r2.font.size = font_size
    r2.font.color.rgb = text_color


def _add_accent_card(slide, accent_text: str, theme: dict):
    """右下角数字/指标卡片"""
    # 卡片背景
    card_bg = slide.shapes.add_shape(1, Inches(10.3), Inches(2.5), Inches(2.7), Inches(3.0))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = theme["accent"]
    card_bg.line.fill.background()

    # 透明度设置（50%）
    try:
        from pptx.oxml.ns import qn
        solidFill = card_bg.fill._xPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                from lxml import etree
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', '40000')  # 40% 不透明
    except Exception:
        pass

    # 数字文字
    num_box = slide.shapes.add_textbox(Inches(10.3), Inches(3.2), Inches(2.7), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = accent_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme["white"]


async def create_pptx_task(task_params: dict[str, Any]) -> dict[str, Any]:
    os.makedirs(STORAGE_DIR, exist_ok=True)

    deliverable_id = task_params.get("deliverableId", str(uuid.uuid4()))
    title = task_params.get("title", "方案演示")
    slides_data: list[dict] = task_params.get("slides", [])
    company_name = task_params.get("companyName", "")
    style = task_params.get("style", "tech")

    theme = _get_theme(style)

    prs = Presentation()
    # 16:9 幻灯片尺寸
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    _add_cover_slide(prs, title, company_name, theme)

    # 2. 内容页
    for slide_info in slides_data:
        s_title = slide_info.get("title", "")

        # ── 核心修复：优先读 bullets 数组，回退到 content 换行文本 ──
        bullets_list = slide_info.get("bullets", [])
        if bullets_list and isinstance(bullets_list, list):
            bullets = [str(b).strip() for b in bullets_list if str(b).strip()]
        else:
            raw_content = slide_info.get("content", "")
            bullets = [line.strip() for line in raw_content.split("\n") if line.strip()]
            if not bullets and raw_content:
                bullets = [raw_content]

        s_notes = slide_info.get("notes", "")
        s_accent = str(slide_info.get("accent", "")).strip()

        _add_content_slide(prs, s_title, bullets, s_notes, s_accent, theme)

    # 3. 结尾页（感谢页）
    end_slide_layout = prs.slide_layouts[6]
    end_slide = prs.slides.add_slide(end_slide_layout)
    _set_slide_bg(end_slide, theme["bg"])

    # 感谢页装饰条
    end_bar = end_slide.shapes.add_shape(1, Inches(0), Inches(3.3), Inches(13.33), Inches(0.9))
    end_bar.fill.solid()
    end_bar.fill.fore_color.rgb = theme["accent"]
    end_bar.line.fill.background()

    end_box = end_slide.shapes.add_textbox(Inches(2), Inches(3.35), Inches(9), Inches(0.85))
    p = end_box.text_frame.paragraphs[0]
    p.text = "感谢聆听  ·  Thank You"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["header_text"] if _is_dark_theme(theme) else theme["white"]
    p.alignment = PP_ALIGN.CENTER

    # 公司名（感谢页副文字）
    if company_name:
        sub_box = end_slide.shapes.add_textbox(Inches(2), Inches(4.4), Inches(9), Inches(0.7))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = company_name
        sp.font.size = Pt(18)
        sp.font.color.rgb = theme["text"] if _is_dark_theme(theme) else theme["subtext"]
        sp.alignment = PP_ALIGN.CENTER

    # 保存文件
    filename = f"{deliverable_id}.pptx"
    filepath = os.path.join(STORAGE_DIR, filename)
    prs.save(filepath)

    return {"fileUrl": f"/files/{filename}"}
